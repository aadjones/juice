# modules/report.py  – v2  (KPI table on slide 2, chart on slide 3)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pathlib import Path


GREEN  = RGBColor(0x0C, 0xAA, 0x41)
YELLOW = RGBColor(0xE9, 0xB8, 0x00)
RED    = RGBColor(0xD4, 0x26, 0x26)


def _kpi_colour(label, value):
    if label == "GQ":
        return GREEN if value >= 1.5 else YELLOW if value >= 1.0 else RED
    if label == "Φf":
        return GREEN if value > 0 else RED
    if label == "Ξσ":
        return GREEN if value < 1.5 else YELLOW if value < 2.5 else RED
    if label == "Fortitude":
        return GREEN if value >= 20 else YELLOW if value >= 10 else RED
    return RGBColor(0, 0, 0)


def build_deck(df, chart_paths, out_dir="report_exports"):
    Path(out_dir).mkdir(exist_ok=True)
    end_date = df.date.max().date()
    out_path = Path(out_dir) / f"Psychic_KPIs_{end_date}.pptx"

    prs = Presentation()

    # -- Slide 1 : Title ----------------------------------------------------
    s0 = prs.slides.add_slide(prs.slide_layouts[0])
    s0.shapes.title.text = "Demby Analytics™"
    d0, d1 = df.date.min().date(), df.date.max().date()
    s0.shapes.placeholders[1].text = f"KPI Board Deck\n{d0} – {d1}"

    # -- Slide 2 : KPI table + commentary -----------------------------------
    s1 = prs.slides.add_slide(prs.slide_layouts[5])
    s1.shapes.title.text = "Key Metrics Snapshot"

    labels = ["GQ", "Φf", "Ξσ", "Fortitude"]
    latest = df.iloc[-1]
    values = [
        round(float(latest.gq), 2),
        round(float(latest.focus_flux), 2),
        round(float(latest.sigma), 2),
        round(float(latest.fortitude), 1),
    ]

    left, top, w, h = Inches(0.4), Inches(1.3), Inches(2.1), Inches(1.1)
    for idx, (lab, val) in enumerate(zip(labels, values)):
        r, c = divmod(idx, 2)
        shp = s1.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left + c * (w + Inches(0.25)),
            top  + r * (h + Inches(0.35)),
            w, h,
        )
        fill = _kpi_colour(lab, val)
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        shp.line.color.rgb = RGBColor(255, 255, 255)

        tf = shp.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = lab
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(255, 255, 255)

        p2 = tf.add_paragraph()
        p2.text = str(val)
        p2.font.size = Pt(22)
        p2.font.color.rgb = RGBColor(255, 255, 255)

    # commentary bullet
    trend = latest.focus_flux - df.focus_flux.iloc[-2] if len(df) > 1 else 0
    comment = (
        f"Flux {'improved' if trend>0 else 'declined'} {trend:+0.2f} since last entry; "
        f"{'Fortitude buffer strong.' if latest.fortitude>=20 else 'Fortitude buffer thinning.'}"
    )
    s1.shapes.add_textbox(Inches(0.4), Inches(3.9), Inches(9), Inches(0.8))\
       .text_frame.text = comment

    # --- Slide 3 : Trend overview chart -------------------------------------
    chart_file = chart_paths.get("juice_anx")
    if chart_file and Path(chart_file).is_file():
        s2 = prs.slides.add_slide(prs.slide_layouts[5])
        s2.shapes.title.text = "Juice vs. Anxiety Trend"

        # add at natural size first
        pic = s2.shapes.add_picture(str(chart_file), Inches(0.4), Inches(1.5))

        # calculate max allowed height (slide height − top margin − bottom margin)
        slide_h     = prs.slide_height           # in EMU units
        top_margin  = Inches(1.5)
        bottom_pad  = Inches(0.2)
        max_h       = slide_h - top_margin - bottom_pad

        if pic.height > max_h:
            # scale down proportionally to fit
            scale = max_h / pic.height
            pic.height = int(pic.height * scale)
            pic.width  = int(pic.width  * scale)
    else:
        # Add a placeholder slide if chart is missing
        s2 = prs.slides.add_slide(prs.slide_layouts[5])
        s2.shapes.title.text = "Juice vs. Anxiety Trend"
        s2.shapes.add_textbox(Inches(0.4), Inches(1.5), Inches(9), Inches(1))\
           .text_frame.text = "(Chart image unavailable)"

    # Always save the presentation, even if chart is missing
    prs.save(out_path)
    return out_path
